"""Generate a compact, visual project-review slide deck (.pptx).

Phase covered: "Design & Development" (15 May 2026 - 14 June 2026), so progress
is framed around the completed implementation baseline: data pipeline, model
training, automated champion selection, drift monitoring, serving, containers,
CI, and benchmark evidence.

Slides:
    1. Title
    2. Objectives & Scope
    3. System Architecture & Design
    4. Methodology Adopted
    5. Core Contribution: Automated Selection & Drift Monitoring
    6. Progress & Key Findings
    7. Proof of Value: Manual vs Automated
    8. Future Work

Usage:
    python report/build_ppt.py
"""
from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

PROJECT_ROOT = Path(__file__).resolve().parents[1]
ARTIFACTS = PROJECT_ROOT / "artifacts"
FIGURES = PROJECT_ROOT / "reports" / "figures"
OUT_PATH = PROJECT_ROOT / "report" / "2024da04367-digit-recognition-review.pptx"

# Palette
NAVY = RGBColor(0x1F, 0x3A, 0x5F)
NAVY_DK = RGBColor(0x15, 0x29, 0x44)
ORANGE = RGBColor(0xDD, 0x84, 0x52)
GREEN = RGBColor(0x55, 0xA8, 0x68)
RED = RGBColor(0xC4, 0x4E, 0x52)
BLUE = RGBColor(0x4C, 0x72, 0xB0)
LIGHT = RGBColor(0xF2, 0xF4, 0xF7)
GREY = RGBColor(0x5A, 0x63, 0x72)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

EMU_W = Inches(13.333)
EMU_H = Inches(7.5)


def load_results():
    defaults = {
        "underfit": {"model_type": "decision_tree", "train_accuracy": 0.323,
                     "val_accuracy": 0.316, "test_accuracy": 0.303,
                     "train_val_gap": 0.007},
        "balanced": {"model_type": "svm", "train_accuracy": 1.0,
                     "val_accuracy": 0.990, "test_accuracy": 0.983,
                     "train_val_gap": 0.010},
        "overfit": {"model_type": "decision_tree", "train_accuracy": 1.0,
                    "val_accuracy": 0.830, "test_accuracy": 0.786,
                    "train_val_gap": 0.170},
    }
    comp_path = ARTIFACTS / "comparison.json"
    if comp_path.exists():
        with open(comp_path, encoding="utf-8") as fh:
            defaults.update(json.load(fh))
    lat = {"p50": 0.237, "p95": 0.398, "p99": 0.718, "rps": 3805}
    lat_path = ARTIFACTS / "latency_benchmark.json"
    if lat_path.exists():
        with open(lat_path, encoding="utf-8") as fh:
            data = json.load(fh)
        for r in data["results"]:
            if r["model"] == "balanced":
                s = r["single"]
                lat = {"p50": s["p50_ms"], "p95": s["p95_ms"],
                       "p99": s["p99_ms"], "rps": s["throughput_rps"]}
    proof = {
        "champion": "balanced",
        "model_type": "svm",
        "selection_ms": 0.05,
        "first_alert": "severity 0.2",
        "errors_caught": 311,
        "dataset": "digits",
        "table": [],
        "manual_acc": None,
        "auto_acc": None,
    }
    champion_path = PROJECT_ROOT / "models" / "champion.json"
    if champion_path.exists():
        with open(champion_path, encoding="utf-8") as fh:
            champion = json.load(fh)
        proof.update({
            "champion": champion.get("champion", proof["champion"]),
            "model_type": champion.get("model_type", proof["model_type"]),
            "selection_ms": champion.get("selection_seconds", 0.00005) * 1000,
            "dataset": champion.get("dataset", proof["dataset"]),
        })
    exp_path = ARTIFACTS / "experiment_results.json"
    if exp_path.exists():
        with open(exp_path, encoding="utf-8") as fh:
            exp = json.load(fh)
        first_alert = exp.get("first_alert") or {}
        if first_alert:
            proof["first_alert"] = f"severity {first_alert.get('severity', 0.2)}"
            proof["errors_caught"] = first_alert.get("errors", proof["errors_caught"])
        proof["dataset"] = exp.get("dataset", proof["dataset"])
        proof["table"] = exp.get("proof_table", [])
        proof["manual_acc"] = exp.get("manual", {}).get("metrics", {}).get(
            "test_accuracy")
        proof["auto_acc"] = exp.get("automated", {}).get("metrics", {}).get(
            "test_accuracy")
    return defaults, lat, proof


# ---------------------------------------------------------------------------
# Drawing helpers
# ---------------------------------------------------------------------------
def _no_line(shape) -> None:
    shape.line.fill.background()


def fill_rect(slide, x, y, w, h, color, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line is None:
        _no_line(shp)
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp


def rounded(slide, x, y, w, h, color, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line is None:
        _no_line(shp)
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1.25)
    shp.shadow.inherit = False
    return shp


def textbox(slide, x, y, w, h, lines, size=14, color=NAVY, bold=False,
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    """lines: str OR list of (text, size, color, bold, bullet, level)."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if isinstance(lines, str):
        lines = [(lines, size, color, bold, False, 0)]
    for i, item in enumerate(lines):
        text, sz, col, bd, bullet, level = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.level = level
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = (("\u2022  " if bullet else "") + text)
        run.font.size = Pt(sz)
        run.font.bold = bd
        run.font.color.rgb = col
        run.font.name = font
    return tb


def title_bar(slide, number, title, subtitle=None):
    fill_rect(slide, 0, 0, EMU_W, Inches(1.15), NAVY)
    fill_rect(slide, 0, Inches(1.15), EMU_W, Inches(0.07), ORANGE)
    # number chip
    chip = rounded(slide, Inches(0.45), Inches(0.27), Inches(0.62), Inches(0.62), ORANGE)
    ctf = chip.text_frame
    ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
    cp = ctf.paragraphs[0]
    cp.alignment = PP_ALIGN.CENTER
    cr = cp.add_run()
    cr.text = str(number)
    cr.font.size = Pt(24)
    cr.font.bold = True
    cr.font.color.rgb = WHITE
    textbox(slide, Inches(1.3), Inches(0.18), Inches(11.5), Inches(0.85),
            [(title, 30, WHITE, True, False, 0)], anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        textbox(slide, Inches(1.32), Inches(0.78), Inches(11.5), Inches(0.35),
                [(subtitle, 13, RGBColor(0xCB, 0xD6, 0xE6), False, False, 0)])


def footer(slide, idx):
    textbox(slide, Inches(0.45), Inches(7.05), Inches(9), Inches(0.35),
            [("Jhanvi Titoria  |  2024da04367  |  M.Tech. DSE  |  BITS Pilani WILP",
              10, GREY, False, False, 0)])
    textbox(slide, Inches(12.3), Inches(7.05), Inches(0.8), Inches(0.35),
            [(str(idx), 10, GREY, True, False, 0)], align=PP_ALIGN.RIGHT)


def add_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank


def pct(x):
    return f"{x * 100:.1f}%"


# ---------------------------------------------------------------------------
# Slides
# ---------------------------------------------------------------------------
def slide_title(prs):
    s = add_slide(prs)
    fill_rect(s, 0, 0, EMU_W, EMU_H, NAVY)
    fill_rect(s, 0, Inches(5.0), EMU_W, Inches(0.08), ORANGE)
    fill_rect(s, 0, Inches(0), Inches(0.28), EMU_H, ORANGE)
    textbox(s, Inches(0.9), Inches(1.0), Inches(11.5), Inches(0.5),
            [("PROJECT REVIEW  \u00b7  PHASE 2: DESIGN & DEVELOPMENT",
              15, ORANGE, True, False, 0)])
    textbox(s, Inches(0.9), Inches(1.7), Inches(11.6), Inches(2.2),
            [("MLOps Pipeline for a", 40, WHITE, True, False, 0),
             ("Digit Recognition Service", 40, WHITE, True, False, 0)])
    textbox(s, Inches(0.92), Inches(3.9), Inches(11.5), Inches(0.9),
            [("Automated Model Selection & Drift Monitoring on an "
              "Experiment-Tracked, Containerised Serving Pipeline", 18,
              RGBColor(0xCB, 0xD6, 0xE6), False, False, 0)])
    textbox(s, Inches(0.9), Inches(5.3), Inches(11.5), Inches(1.5),
            [("Jhanvi Titoria", 20, WHITE, True, False, 0),
             ("BITS ID: 2024da04367  \u00b7  M.Tech. Data Science & Engineering",
              14, RGBColor(0xCB, 0xD6, 0xE6), False, False, 0),
             ("Supervisor: Abhishek Kumar  \u00b7  Examiner: Abhishek Kumar Jain",
              14, RGBColor(0xCB, 0xD6, 0xE6), False, False, 0),
             ("Reporting window: 15 May 2026 \u2013 14 June 2026", 14,
              ORANGE, True, False, 0)])


def card(slide, x, y, w, h, heading, hcolor, items):
    rounded(slide, x, y, w, h, LIGHT)
    fill_rect(slide, x, y, w, Inches(0.55), hcolor)
    textbox(slide, x + Inches(0.2), y + Inches(0.05), w - Inches(0.3), Inches(0.5),
            [(heading, 16, WHITE, True, False, 0)], anchor=MSO_ANCHOR.MIDDLE)
    body = [(t, 13, NAVY, False, True, 0) for t in items]
    textbox(slide, x + Inches(0.25), y + Inches(0.7), w - Inches(0.45),
            h - Inches(0.85), body, anchor=MSO_ANCHOR.TOP)


def slide_objectives(prs):
    s = add_slide(prs)
    title_bar(s, 1, "Objectives & Scope",
              "What the project sets out to build and where its boundaries lie")
    card(s, Inches(0.45), Inches(1.5), Inches(6.1), Inches(5.2),
         "OBJECTIVES", BLUE, [
             "Build an end-to-end, CPU-only MLOps pipeline for handwritten "
             "digit recognition.",
             "Train three contrasted models (underfit / balanced / overfit) "
             "to study generalisation.",
             "Track every run in MLflow: params, metrics, train\u2013val gap, "
             "model artifacts.",
             "Promote a champion automatically using a transparent "
             "generalisation-aware rule.",
             "Monitor live labelled batches and raise drift / performance-drop "
             "alerts.",
             "Serve the model via FastAPI with /predict, /livez, /readyz, "
             "/monitor and /metrics.",
             "Containerise training & serving; add tests, CI and latency/load "
             "benchmarks.",
         ])
    card(s, Inches(6.8), Inches(1.5), Inches(6.1), Inches(5.2),
         "SCOPE", GREEN, [
             "Dataset: scikit-learn 8\u00d78 digits \u2014 1,797 samples, 64 "
             "features, 10 classes; MNIST switch available through config.",
             "Reproducible stratified train / validation / test splits.",
             "Three-model experimental arc with automated under/overfitting "
             "diagnosis.",
             "Serving, Docker, automated tests, CI, latency characterisation "
             "and Prometheus observability.",
             "In scope: operational engineering, model governance and model "
             "monitoring.",
             "Out of scope: large-scale GPU training and novel model "
             "research.",
         ])
    footer(s, 2)


def chevron(slide, x, y, w, h, color, label, sub):
    shp = slide.shapes.add_shape(MSO_SHAPE.PENTAGON, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    _no_line(shp)
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.size = Pt(13)
    r.font.bold = True
    r.font.color.rgb = WHITE
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = sub
    r2.font.size = Pt(9)
    r2.font.color.rgb = RGBColor(0xEC, 0xF1, 0xF7)


def slide_methodology(prs):
    s = add_slide(prs)
    title_bar(s, 3, "Methodology Adopted",
              "A reproducible, config-driven MLOps lifecycle implemented in Phase 2")
    steps = [
        (NAVY, "DATA", "stratified splits"),
        (BLUE, "TRAIN", "YAML configs"),
        (RGBColor(0x6A, 0x51, 0xA3), "TRACK", "MLflow"),
        (ORANGE, "SELECT", "champion rule"),
        (GREEN, "SERVE", "FastAPI + metrics"),
        (RED, "MONITOR", "drift alerts"),
    ]
    n = len(steps)
    x0 = Inches(0.45)
    gap = Inches(0.05)
    total = EMU_W - Inches(0.9)
    w = Emu(int((total - gap * (n - 1)) / n)) + Inches(0.35)
    for i, (color, label, sub) in enumerate(steps):
        x = Emu(int(x0) + i * int(Emu(int((total - gap * (n - 1)) / n)) + gap))
        chevron(s, x, Inches(1.6), w, Inches(0.95), color, label, sub)

    # Detail cards
    card(s, Inches(0.45), Inches(2.9), Inches(4.0), Inches(3.6),
         "DESIGN PRINCIPLES", BLUE, [
             "Config-driven runs (configs/*.yaml) \u2014 reproducible & "
             "auditable.",
             "Single random seed across all splits.",
             "Scaler + estimator wrapped in one sklearn Pipeline.",
             "Same data feeds all three models for a fair comparison.",
         ])
    card(s, Inches(4.65), Inches(2.9), Inches(4.0), Inches(3.6),
         "THREE-MODEL ARC", ORANGE, [
             "Underfit: depth-2 decision tree (high bias).",
             "Balanced: RBF-kernel SVM (good fit).",
             "Overfit: fully grown decision tree (high variance).",
             "Auto-verdict from train\u2013validation gap.",
             "Champion rule excludes overfit candidates by max-gap.",
         ])
    card(s, Inches(8.85), Inches(2.9), Inches(4.05), Inches(3.6),
         "TOOLING", GREEN, [
             "scikit-learn, MLflow, FastAPI, Pydantic.",
             "Prometheus client for /metrics.",
             "Docker train / serve images.",
             "pytest, Locust, GitHub Actions CI.",
         ])
    footer(s, 4)


def metric_chip(slide, x, y, w, value, label, color):
    rounded(slide, x, y, w, Inches(1.15), color)
    textbox(slide, x, y + Inches(0.12), w, Inches(0.6),
            [(value, 26, WHITE, True, False, 0)], align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE)
    textbox(slide, x, y + Inches(0.72), w, Inches(0.4),
            [(label, 11, WHITE, False, False, 0)], align=PP_ALIGN.CENTER)


def results_table(slide, x, y, comp):
    rows = [
        ("Run", "Model", "Train", "Val", "Test", "Gap", "Verdict"),
        ("Underfit", comp["underfit"]["model_type"],
         pct(comp["underfit"]["train_accuracy"]),
         pct(comp["underfit"]["val_accuracy"]),
         pct(comp["underfit"]["test_accuracy"]),
         f"{comp['underfit']['train_val_gap']*100:+.1f}", "Underfit"),
        ("Balanced", comp["balanced"]["model_type"],
         pct(comp["balanced"]["train_accuracy"]),
         pct(comp["balanced"]["val_accuracy"]),
         pct(comp["balanced"]["test_accuracy"]),
         f"{comp['balanced']['train_val_gap']*100:+.1f}", "Good fit"),
        ("Overfit", comp["overfit"]["model_type"],
         pct(comp["overfit"]["train_accuracy"]),
         pct(comp["overfit"]["val_accuracy"]),
         pct(comp["overfit"]["test_accuracy"]),
         f"{comp['overfit']['train_val_gap']*100:+.1f}", "Overfit"),
    ]
    n_rows, n_cols = len(rows), len(rows[0])
    gtbl = slide.shapes.add_table(n_rows, n_cols, x, y, Inches(7.7),
                                  Inches(2.0)).table
    widths = [1.0, 1.4, 0.9, 0.9, 0.9, 0.8, 1.1]
    for j, wd in enumerate(widths):
        gtbl.columns[j].width = Inches(wd)
    verdict_colors = {"Underfit": RED, "Good fit": GREEN, "Overfit": RED}
    for i, row in enumerate(rows):
        gtbl.rows[i].height = Inches(0.42)
        for j, val in enumerate(row):
            cell = gtbl.cell(i, j)
            cell.margin_top = Pt(2)
            cell.margin_bottom = Pt(2)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = val
            run.font.size = Pt(11 if i else 12)
            run.font.name = "Calibri"
            if i == 0:
                run.font.bold = True
                run.font.color.rgb = WHITE
                cell.fill.solid()
                cell.fill.fore_color.rgb = NAVY
            else:
                run.font.color.rgb = NAVY
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if i % 2 else LIGHT
                if j == 6:
                    run.font.bold = True
                    run.font.color.rgb = verdict_colors.get(val, NAVY)


def slide_progress(prs, comp, lat, proof):
    s = add_slide(prs)
    title_bar(s, 5, "Progress & Key Findings",
              "Status at end of Design & Development (15 May\u201314 Jun 2026)")
    # progress chips row
    metric_chip(s, Inches(0.45), Inches(1.45), Inches(3.0),
                "100%", "Design & development complete", GREEN)
    metric_chip(s, Inches(3.6), Inches(1.45), Inches(3.0),
                pct(comp["balanced"]["test_accuracy"]),
                "Best model test accuracy", BLUE)
    metric_chip(s, Inches(6.75), Inches(1.45), Inches(3.0),
                f"{comp['overfit']['train_val_gap']*100:.0f} pp",
                "Overfit train\u2013val gap", RED)
    metric_chip(s, Inches(9.9), Inches(1.45), Inches(3.0),
                f"{proof['selection_ms']:.2f} ms", "Champion selection time", ORANGE)

    textbox(s, Inches(0.45), Inches(2.85), Inches(7.7), Inches(0.4),
            [("Implemented three-model comparison and champion baseline",
              15, NAVY, True, False, 0)])
    results_table(s, Inches(0.45), Inches(3.3), comp)

    card(s, Inches(8.45), Inches(2.85), Inches(4.45), Inches(3.85),
         "KEY FINDINGS", NAVY, [
             "Data, training, analysis, selector, monitor and serving modules "
             "are implemented.",
             "Bias\u2013variance spectrum reproduced across underfit / balanced / "
             "overfit runs.",
             f"Balanced RBF-SVM: {pct(comp['balanced']['val_accuracy'])} val, "
             f"only {comp['balanced']['train_val_gap']*100:.1f} pp gap.",
             f"Overfit tree: 100% train but {pct(comp['overfit']['val_accuracy'])} "
             "val \u2014 clear over-fitting.",
             f"Automated selector promotes {proof['champion']} "
             f"({proof['model_type']}) as champion.",
             f"Monitor alerts at {proof['first_alert']} and flags "
             f"{proof['errors_caught']} errors before users.",
         ])

    fig = FIGURES / "accuracy_comparison.png"
    if fig.exists():
        s.shapes.add_picture(str(fig), Inches(0.55), Inches(5.45),
                             height=Inches(1.25))
    footer(s, 6)


def slide_future(prs):
    s = add_slide(prs)
    title_bar(s, 8, "Future Work",
              "Planned activities for the upcoming phases")
    timeline = [
        (GREEN, "Outline", "01 May \u2013 14 May 2026",
         ["Literature review and problem framing",
          "Project outline and first prototype completed"]),
        (GREEN, "Design & Development", "15 May \u2013 14 Jun 2026",
         ["Data, training, selector, monitor and service modules completed",
          "Dockerfiles, CI workflow and report generators added"]),
        (ORANGE, "Testing", "15 Jun \u2013 01 Jul 2026",
         ["Expand pytest + GitHub Actions CI",
          "Validate /predict, /monitor and /metrics endpoints",
          "Latency benchmark + Locust load tests"]),
        (NAVY, "Submission", "13 Jul \u2013 17 Jul 2026",
         ["Supervisor / examiner review",
          "Final dissertation packaging and submission"]),
    ]
    n = len(timeline)
    x0 = 0.45
    cw = 3.02
    gapw = 0.16
    for i, (color, phase, dates, items) in enumerate(timeline):
        x = Inches(x0 + i * (cw + gapw))
        rounded(s, x, Inches(1.6), Inches(cw), Inches(0.95), color)
        textbox(s, x + Inches(0.12), Inches(1.66), Inches(cw - 0.2), Inches(0.5),
                [(phase, 14, WHITE, True, False, 0)], anchor=MSO_ANCHOR.MIDDLE)
        textbox(s, x + Inches(0.12), Inches(2.12), Inches(cw - 0.2), Inches(0.35),
                [(dates, 10, RGBColor(0xEC, 0xF1, 0xF7), False, False, 0)])
        rounded(s, x, Inches(2.7), Inches(cw), Inches(3.4), LIGHT)
        body = [(t, 12, NAVY, False, True, 0) for t in items]
        textbox(s, x + Inches(0.18), Inches(2.85), Inches(cw - 0.32),
                Inches(3.1), body)
    textbox(s, Inches(0.45), Inches(6.35), Inches(12.4), Inches(0.5),
            [("Risk mitigation: CPU-only design keeps the pipeline portable; "
              "digits/MNIST dataset switch keeps experiments reproducible while "
              "allowing scale-up.",
              12, GREY, False, False, 0)])
    footer(s, 8)


PURPLE = RGBColor(0x6A, 0x51, 0xA3)


def flow_box(slide, x, y, w, h, title, sub, color):
    rounded(slide, x, y, w, h, color)
    textbox(slide, x, y + Inches(0.14), w, Inches(0.45),
            [(title, 13, WHITE, True, False, 0)], align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE)
    textbox(slide, x, y + Inches(0.6), w, Inches(0.35),
            [(sub, 9, RGBColor(0xEC, 0xF1, 0xF7), False, False, 0)],
            align=PP_ALIGN.CENTER)


def arrow(slide, x, y, w, h, color):
    shp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    _no_line(shp)
    shp.shadow.inherit = False


def slide_architecture(prs):
    s = add_slide(prs)
    title_bar(s, 2, "System Architecture & Design",
              "End-to-end data and model flow built in the Design & Development phase")
    stages = [
        (NAVY, "DATA", "digits / MNIST"),
        (BLUE, "TRAIN", "YAML configs"),
        (PURPLE, "MLflow", "tracking"),
        (ORANGE, "SELECTOR", "champion rule"),
        (GREEN, "SERVE", "FastAPI"),
        (RED, "MONITOR", "drift alerts"),
    ]
    bx, by, bw, bh = 0.5, 1.45, 1.78, 1.1
    aw = 0.3
    step = bw + aw
    for i, (color, label, sub) in enumerate(stages):
        x = bx + i * step
        flow_box(s, Inches(x), Inches(by), Inches(bw), Inches(bh), label, sub, color)
        if i < len(stages) - 1:
            arrow(s, Inches(x + bw + 0.02), Inches(by + 0.42),
                  Inches(aw - 0.04), Inches(0.26), GREY)

    textbox(s, Inches(0.5), Inches(2.7), Inches(12.4), Inches(0.4),
            [("Artifacts flow:  models/champion.joblib + champion.json   \u00b7   "
              "artifacts/metrics/*.json   \u00b7   reports/figures/*", 12,
              GREY, False, False, 0)], align=PP_ALIGN.CENTER)

    card(s, Inches(0.45), Inches(3.2), Inches(6.1), Inches(2.95),
         "CROSS-CUTTING ENGINEERING", BLUE, [
             "Containerisation: separate Docker train / serve images.",
             "CI: GitHub Actions runs the pytest suite on every push.",
             "Quality: pytest across data, models, serving, selector, monitor.",
             "Performance: inference latency benchmark + Locust load test.",
         ])
    card(s, Inches(6.8), Inches(3.2), Inches(6.1), Inches(2.95),
         "OBSERVABILITY & OPS", GREEN, [
             "Prometheus /metrics: predictions, latency histogram, accuracy "
             "gauge, drift-alert counter.",
             "Health probes /livez, /readyz \u2014 Kubernetes-ready.",
             "Config-driven dataset switch (digits \u2194 MNIST) with subsampling.",
             "Champion model auto-loaded by the service at startup.",
         ])
    textbox(s, Inches(0.45), Inches(6.3), Inches(12.4), Inches(0.55),
            [("Design decisions: scaler + estimator in one sklearn Pipeline  \u00b7  "
              "single random seed for reproducibility  \u00b7  gap-based champion "
              "rule  \u00b7  file-based champion for offline, auditable promotion.",
              12, GREY, False, False, 0)])
    footer(s, 3)


def slide_contribution(prs, proof):
    s = add_slide(prs)
    title_bar(s, 4, "Core Contribution \u2014 Automated Selection & Drift Monitoring",
              "The original module that replaces manual model picking and watching")
    rounded(s, Inches(0.45), Inches(1.4), Inches(12.45), Inches(0.72), NAVY_DK)
    textbox(s, Inches(0.7), Inches(1.4), Inches(12.0), Inches(0.72),
            [("Research question: Can an automated selection-and-monitoring "
              "module choose better-generalising models and catch performance "
              "degradation faster than manual selection?", 14, WHITE, True,
              False, 0)], anchor=MSO_ANCHOR.MIDDLE)

    card(s, Inches(0.45), Inches(2.35), Inches(6.1), Inches(3.2),
         "SELECTOR \u2014 champion picker", ORANGE, [
             "Reads every run's metrics from artifacts/metrics/.",
             "Rule: highest validation accuracy with train\u2013val gap \u2264 0.08 "
             "(over-fit runs excluded).",
             "Promotes models/champion.joblib + champion.json and records the "
             "baseline accuracy.",
             "Replaces hard-coded model choice \u2014 auditable & reproducible.",
         ])
    card(s, Inches(6.8), Inches(2.35), Inches(6.1), Inches(3.2),
         "MONITOR \u2014 drift / drop alerts", RED, [
             "Scores the champion on fresh labelled batches.",
             "Alerts when accuracy falls more than a threshold below baseline.",
             "Live via POST /monitor; Prometheus accuracy gauge + alert counter.",
             "Non-zero exit code \u2192 CI / Alertmanager hook for auto-response.",
         ])

    chip_w = 3.0
    metric_chip(s, Inches(0.45), Inches(5.7), Inches(chip_w),
                f"{proof['champion']} / {proof['model_type']}",
                "Auto-selected champion", ORANGE)
    metric_chip(s, Inches(3.6), Inches(5.7), Inches(chip_w),
                f"{proof['selection_ms']:.2f} ms", "Selection time", BLUE)
    metric_chip(s, Inches(6.75), Inches(5.7), Inches(chip_w),
                proof["first_alert"], "First drift alert", RED)
    metric_chip(s, Inches(9.9), Inches(5.7), Inches(chip_w),
                str(proof["errors_caught"]), "Errors caught before users", GREEN)
    footer(s, 5)


def proof_table_native(slide, x, y, rows):
    data = [("Metric", "WITHOUT module (manual)", "WITH module (automated)")]
    for r in rows:
        data.append((r["metric"], r["without_module"], r["with_module"]))
    n_rows, n_cols = len(data), 3
    tbl = slide.shapes.add_table(n_rows, n_cols, x, y, Inches(12.45),
                                 Inches(0.42 * n_rows)).table
    for j, wd in enumerate((3.3, 4.55, 4.6)):
        tbl.columns[j].width = Inches(wd)
    for i, row in enumerate(data):
        tbl.rows[i].height = Inches(0.42)
        for j, val in enumerate(row):
            cell = tbl.cell(i, j)
            cell.margin_top = Pt(2)
            cell.margin_bottom = Pt(2)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            run = p.add_run()
            run.text = val
            run.font.size = Pt(11 if i else 12)
            run.font.name = "Calibri"
            cell.fill.solid()
            if i == 0:
                run.font.bold = True
                run.font.color.rgb = WHITE
                cell.fill.fore_color.rgb = NAVY
            else:
                cell.fill.fore_color.rgb = WHITE if i % 2 else LIGHT
                run.font.color.rgb = NAVY
                if j == 2:
                    run.font.bold = True
                    run.font.color.rgb = GREEN


def slide_proof(prs, proof):
    s = add_slide(prs)
    title_bar(s, 6, "Proof of Value \u2014 Manual vs Automated",
              f"Same data & candidates, run two ways (dataset: {proof['dataset']})")
    rows = proof.get("table") or [
        {"metric": "Model served", "without_module": "balanced (svm) \u2014 manual",
         "with_module": "balanced (svm) \u2014 auto"},
        {"metric": "Drift detection", "without_module": "none",
         "with_module": "alerts at severity 0.2"},
        {"metric": "Errors caught before users", "without_module": "0",
         "with_module": "311 at first alert"},
    ]
    proof_table_native(s, Inches(0.45), Inches(1.45), rows)

    img_top = 1.55 + 0.42 * (len(rows) + 1)
    fig = FIGURES / "drift_sweep.png"
    if fig.exists():
        s.shapes.add_picture(str(fig), Inches(0.55), Inches(img_top),
                             height=Inches(2.35))
    card(s, Inches(6.0), Inches(img_top), Inches(6.9), Inches(2.35),
         "WHAT THIS PROVES", NAVY, [
             "Automated selection matches the manual pick at zero human effort "
             "and is fully reproducible.",
             "The monitor turns silent failure into an explicit, early alert.",
             f"At the first alert it flagged {proof['errors_caught']} errors that "
             "the manual setup would have served unnoticed.",
             "Run --dataset mnist to test whether the best model changes at scale.",
         ])
    footer(s, 7)


def build():
    prs = Presentation()
    prs.slide_width = EMU_W
    prs.slide_height = EMU_H
    comp, lat, proof = load_results()

    slide_title(prs)
    slide_objectives(prs)
    slide_architecture(prs)
    slide_methodology(prs)
    slide_contribution(prs, proof)
    slide_progress(prs, comp, lat, proof)
    slide_proof(prs, proof)
    slide_future(prs)

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PATH)
    print(f"Deck written to: {OUT_PATH}  ({len(prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    build()
